import argparse
import openai
import os
import zipfile
from pptx import Presentation

openai.api_key = os.environ["OPENAI_API_KEY"]

parser = argparse.ArgumentParser()
parser.add_argument("-f", "--file", required=True, help="path to input PPTX file")
args = parser.parse_args()

input_file = args.file
output_file = os.path.splitext(input_file)[0] + "_new.pptx"

# 建立Presentation物件
prs =Presentation(args.file)

for slide in prs.slides:
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                text = run.text
                try:
                    completion = openai.Completion.create(
                        engine="text-davinci-003",
                        prompt=f"Translate the following text to traditional chinese: {text}",
                        max_tokens=1024,
                        n=1,
                        stop=None,
                        temperature=0.5,
                        api_key=openai.api_key,
                    )
                    translated_text = completion.choices[0].text
                    run.text = translated_text
                except Exception as e:
                    print(f"Error: {e}")

# 儲存翻譯好的ppt檔案
prs.save(output_file)
print("己完成!")

